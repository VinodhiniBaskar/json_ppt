# -*- coding: utf-8 -*-
# import tracemalloc
import os
import logging
import uuid
import datetime
import json
import sys
import getopt
import time
import traceback
import random
import re
import timeit
import pickle
import csv
import subprocess
import pafy
from pytube import YouTube
from pydub import AudioSegment
from mhyt import yt_download
import shutil
from functools import wraps
from flask import abort
import marshmallow as ma
from urllib.error import ContentTooShortError
from marshmallow import Schema, post_load, validate, ValidationError
from mv_musictool.mvmodels.Projects import Project,ProjectFile,TempFileStorage
from mv_musictool.mvexception.exception import MVException, ValidationException,Test
from mongoengine.queryset.visitor import Q
from flask import Flask,jsonify
from flask_pymongo import PyMongo
from flask_restplus import Api, Resource, fields
from bson import ObjectId
from mv_musictool import settings
from mv_musictool.api import utils
from pymongo import MongoClient
import threading
from werkzeug.utils import secure_filename
import pptx
import copy
import six
from pptx import Presentation
# from pptx.enum.shapes import MSO_SHAPE_TYPE
# from pptx.util import Inches
from pptx.opc.package import PartFactory
from pptx.parts.media import MediaPart
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from urllib.parse import urljoin
from urllib.request import pathname2url
from pptx.enum.action import PP_ACTION
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
# PartFactory.part_type_for.update(
#     {
#         'audio/mp3': MediaPart
#     }
# )
print (sys.getdefaultencoding())

log = logging.getLogger(__name__)
""" Db initialization """
local=MongoClient()
db=local['test_youtube_db']
coll=db["test_youtube_account"]


WATCH_URL = "https://www.youtube.com/watch?v="


DEFAULT_UPLOAD_PATH = settings.MEDIA_PATH+"/"
BASE_MEDIA_PATH = "mv_musictool/static/media/"
THUMBNAIL_JPG = "_thumbnail.jpg"

#TODO handle mv exception
class ProjectSchema(Schema):
    db_id = ma.fields.Str(allow_none=True)
    
    @post_load
    def make_project(self, data, **kwargs):
        return Project(**data)

'''
    ==============================================
    musictool - Class Factory
    ================================================
'''

class MusicSetFactory(object):

    # db connect in __init__?
    def __init__(self):
        # log.debug ('init')
        pass

    def get_marshalled_schema(self,obj):
        if obj:
            schema=ProjectSchema()
            retdata = schema.dump(obj)
            return retdata
    
    def create_project(self,data):
        i = 0
        proj_obj ={}
        res =[]
        with open('/home/desktop-obs-59/Projects/json_ppt/mv_musictool/input_data.json') as f:
            data = json.load(f)
            
        with open('/home/desktop-obs-59/Projects/json_ppt/mv_musictool/input_json_slide1.json') as f:
            data_first = json.load(f)

        prs = Presentation('/home/desktop-obs-59/Projects/json_ppt/TemplateForJson.pptx')

        
        def path2url(path):
            return urljoin('file:', pathname2url(path))
       

        def add_slide(prs, layout, title, img_path,first_audio_path,second_audio_path,score):
            title_fullname =  "COMPOSITION 1 –" + title
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title
            
            img_left = Inches(0.78)
            img_top = Inches(1.81)
            img_width = Inches(9.24)
            img_height = Inches(3.45)
            add_picture = slide.shapes.add_picture(img_path, img_left, img_top, img_width, img_height)

            shape_left = Cm(27.7)
            shape_top = Cm(6.93)
            shape_width = Cm(2.67)
            shape_height = Cm(2.80)
            add_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, shape_left, shape_top, shape_width, shape_height)
            tf = add_shape.text_frame
            p = tf.add_paragraph()
            r = p.add_run()
            r.font.bold = True
            r.font.size = Pt(14)
            r.font.color.rgb = RGBColor(0, 0, 0)
            r.text = 'Score '+ score

            fill = add_shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(221, 221, 221)
            line = add_shape.line
            line.color.rgb = RGBColor(221, 221, 221)
            line.color.brightness = 0.5
            line.width = Pt(2.5)
            
            # left = Inches(11.0)
            # top = Inches(4)
            # width = Inches(0.88)
            # height = Inches(0.88)
            # music_slide = slide.shapes.add_movie(audio_path,left = Inches(1), top= Inches(1), \
            #                 width= Inches(1), height= Inches(1))
            # path = audio_path
            first_audio_link = path2url(first_audio_path)
            print(first_audio_link)

            # First Audio
            shape_left = Cm(28.02)
            shape_top = Cm(10.43)
            shape_width = Cm(2.25)
            shape_height = Cm(2.25)
            add_shape = slide.shapes.add_shape(MSO_SHAPE.ACTION_BUTTON_SOUND, shape_left, shape_top, shape_width, shape_height)
            fill = add_shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(221, 221, 221)
            line = add_shape.line
            line.color.rgb = RGBColor(221, 221, 221)
            line.color.brightness = 0.5
            line.width = Pt(2.5)
            text_frame = add_shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            # run = p.add_run()
            # run.text = 'foobar'
            # p = shape_.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = 'music'
            hlink = r.hyperlink
            hlink.address = first_audio_link

            # Second Audio
            second_audio_link = path2url(second_audio_path)
            print(second_audio_path)

            shape_left = Cm(28.02)
            shape_top = Cm(13.19)
            shape_width = Cm(2.25)
            shape_height = Cm(2.25)
            add_shape = slide.shapes.add_shape(MSO_SHAPE.ACTION_BUTTON_SOUND, shape_left, shape_top, shape_width, shape_height)
            fill = add_shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(221, 221, 221)
            line = add_shape.line
            line.color.rgb = RGBColor(221, 221, 221)
            line.color.brightness = 0.5
            line.width = Pt(2.5)
            text_frame = add_shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            # run = p.add_run()
            # run.text = 'foobar'
            # p = shape_.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = 'music + vo'
            hlink = r.hyperlink
            hlink.address = second_audio_link

            return slide

        k = 0
        slide = prs.slides[0] #first slide
        for val_1 in data_first:
            # print(" ".join(data_first[val_1]))
            table = slide.shapes[0].table.cell(0,k).text = val_1
            table = slide.shapes[0].table.cell(1,k).text = " ".join(data_first[val_1])
            k += 1
        

        def add_shape_third_slide(column,left,top,width,height,audio_path):
            second_audio_link = path2url(audio_path)
            add_shape = slide.shapes.add_shape(MSO_SHAPE.ACTION_BUTTON_SOUND, left, top, width, height)
            fill = add_shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(221, 221, 221)
            line = add_shape.line
            line.color.rgb = RGBColor(221, 221, 221)
            line.color.brightness = 0.5
            line.width = Pt(2.5)
            text_frame = add_shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            r = p.add_run()
            r.text = 'music'
            hlink = r.hyperlink
            hlink.address = second_audio_link
            return add_shape

        title_slide_layout = prs.slide_layouts[4]
        j =0
        slide = prs.slides[2] #third slide
        for val in data:
            shape_top = Cm(11.00 )
            shape_width = Cm(2.25)
            shape_height = Cm(2.25)
            
            if j == 0:
                shape_left = Cm(2.55)
                print("filepfhjds",data[val]['mp3_first'])
                add_shape_third_slide(j,shape_left,shape_top,shape_width,shape_height,data[val]['mp3_first'])
            if j == 1:
                shape_left = Cm(5.60)
                print("filepfhjds",data[val]['mp3_first'])
                add_shape_third_slide(j,shape_left,shape_top,shape_width,shape_height,data[val]['mp3_first'])
            if j ==2:
                shape_left = Cm(8.40)
                add_shape_third_slide(j,shape_left,shape_top,shape_width,shape_height,data[val]['mp3_first'])
            if j ==3:
                shape_left = Cm(11.55)
                add_shape_third_slide(j,shape_left,shape_top,shape_width,shape_height,data[val]['mp3_first'])
            if j ==4:
                shape_left = Cm(14.60)
                add_shape_third_slide(j,shape_left,shape_top,shape_width,shape_height,data[val]['mp3_first'])
            if j ==5:
                shape_left = Cm(17.55)
                add_shape_third_slide(j,shape_left,shape_top,shape_width,shape_height,data[val]['mp3_first'])
            if j ==6:
                shape_left = Cm(20.55)
                add_shape_third_slide(j,shape_left,shape_top,shape_width,shape_height,data[val]['mp3_first'])
            if j ==7:
                shape_left = Cm(23.80)
                add_shape_third_slide(j,shape_left,shape_top,shape_width,shape_height,data[val]['mp3_first'])
            if j ==8:
                shape_left = Cm(26.60)
                add_shape_third_slide(j,shape_left,shape_top,shape_width,shape_height,data[val]['mp3_first'])
            if j ==9:
                shape_left = Cm(29.55)
                add_shape_third_slide(j,shape_left,shape_top,shape_width,shape_height,data[val]['mp3_first'])
            
            first_audio_link = path2url(data[val]['mp3_first'])
            table = slide.shapes[2].table.cell(0,j).text= val
            table = slide.shapes[2].table.cell(1,j).text= data[val]['score']
            # print(j)
            j += 1
        
        for multiple in data:
            title_fullname = 'COMPOSITION 1 – ' + multiple
            img_path = data[multiple]['img_path']
            score = data[multiple]['score']
            first_audio_path = data[multiple]['mp3_first']
            second_audio_path = data[multiple]['mp3_second']
            slide = add_slide(prs, title_slide_layout,title_fullname, img_path, first_audio_path,second_audio_path,score)
        

        
        # slide = prs.slides[2]
        # audio_path = '/home/desktop-obs-59/Projects/json_ppt/333845 Ireland - Online.mp4'
        # music_slide = slide.shapes.add_movie(audio_path,left = Inches(1), top= Inches(1), \
        #                     width= Inches(1), height= Inches(1))
        # add_picture = slide.shapes.add_picture(img_path, left = Inches(1), top= Inches(1), \
        #                     width= Inches(1), height= Inches(1))
        prs.save("PPT_final.pptx")

        # outfile_path = '/json_ppt/out.pptx'
        return ''
                   
musicSetFactory = MusicSetFactory()



